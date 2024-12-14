from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# Function to style text
def style_text(shape, font_name="Open Sans", font_size=24, font_color=(0, 0, 0)):
    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.name = font_name
            run.font.size = Pt(font_size)
            run.font.color.rgb = RGBColor(*font_color)

# Function to set slide background color
def set_slide_background_color(slide, rgb_color=(255, 255, 255)):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(*rgb_color)

# # Function to add a rectangle shape
# def add_rectangle(slide, left, top, width, height, fill_color=(91, 155, 213)):
#     shape = slide.shapes.add_shape(
#         MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height)
#     )
#     shape.fill.solid()
#     shape.fill.fore_color.rgb = RGBColor(*fill_color)
#     shape.line.color.rgb = RGBColor(0, 0, 0)  # Border color

# # Function to add a logo
# def add_logo(slide, image_path, left=0, top=0, width=1.0, height=1.0):
#     slide.shapes.add_picture(image_path, Inches(left), Inches(top), Inches(width), Inches(height))

# Function to add a slide with styled content
def add_slide_with_styling(presentation, title, content):
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])  # Title and Content layout
    slide.shapes.title.text = title
    slide.placeholders[1].text = content

    # Apply styles
    style_text(slide.shapes.title, font_name="Calibri", font_size=36, font_color=(0, 51, 102))  # Title
    style_text(slide.placeholders[1], font_name="Calibri", font_size=20, font_color=(60, 60, 60))  # Content
    set_slide_background_color(slide, rgb_color=(240, 240, 240))  # Light gray background

    return slide

# Main script to create the presentation
def create_presentation():
    presentation = Presentation()

    # Example slide content
    slides_content = [
        ("Welcome to Create a Presentation with Python", 
         "Learn how to create a PowerPoint presentation with Python using the python-pptx library."),
        ("Why Use Python for Presentations?", 
         "- Easy to learn and use\n- Cross-platform\n- Can be automated"),
        ("How to Use Python for Presentations",
         "- Install python-pptx library\n- Write Python code to create slides\n- Run the code"),
        ("Example: Creating a Presentation",
         "- Create a new presentation\n- Add slides with content\n- Customize styling\n- Save the presentation"),
        ("Conclusion",
         "- python-pptx is powerful for presentation automation")
    ]

    for title, content in slides_content:
        slide = add_slide_with_styling(presentation, title, content)

        # Add a rectangle for decoration
        # add_rectangle(slide, 0.5, 1.0, 8.0, 0.5, fill_color=(91, 155, 213))

    # Save the presentation
    presentation.save("./My_Deck.pptx")
    print("Presentation saved as My_Deck.pptx")

# Run the script
if __name__ == "__main__":
    create_presentation()
