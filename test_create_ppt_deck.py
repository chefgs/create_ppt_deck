import pytest
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
from create_ppt_deck import style_text, set_slide_background_color, add_slide_with_styling

def test_style_text():
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    title_shape = slide.shapes.title
    title_shape.text = "Test Title"

    style_text(title_shape, font_name="Arial", font_size=30, font_color=(255, 0, 0))

    for paragraph in title_shape.text_frame.paragraphs:
        for run in paragraph.runs:
            assert run.font.name == "Arial"
            assert run.font.size == Pt(30)
            assert run.font.color.rgb == RGBColor(255, 0, 0)

def test_set_slide_background_color():
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])

    set_slide_background_color(slide, rgb_color=(0, 255, 0))

    background = slide.background
    fill = background.fill
    assert fill.fore_color.rgb == RGBColor(0, 255, 0)

def test_add_slide_with_styling():
    presentation = Presentation()
    title = "Test Slide"
    content = "This is a test slide content."

    slide = add_slide_with_styling(presentation, title, content)

    assert slide.shapes.title.text == title
    assert slide.placeholders[1].text == content

    # Check title style
    for paragraph in slide.shapes.title.text_frame.paragraphs:
        for run in paragraph.runs:
            assert run.font.name == "Calibri"
            assert run.font.size == Pt(36)
            assert run.font.color.rgb == RGBColor(0, 51, 102)

    # Check content style
    for paragraph in slide.placeholders[1].text_frame.paragraphs:
        for run in paragraph.runs:
            assert run.font.name == "Calibri"
            assert run.font.size == Pt(20)
            assert run.font.color.rgb == RGBColor(60, 60, 60)

    # Check background color
    background = slide.background
    fill = background.fill
    assert fill.fore_color.rgb == RGBColor(240, 240, 240)

if __name__ == "__main__":
    pytest.main()